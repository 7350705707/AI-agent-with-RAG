"""Document loaders — extract text from PDF, DOCX, and PPTX files."""

import logging
from pathlib import Path

from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
)
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

from app.config import CHUNK_SIZE, CHUNK_OVERLAP

log = logging.getLogger(__name__)

SPLITTER = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)


def _load_pptx(file_path: str) -> list[Document]:
    """Load a PPTX file using python-pptx (fast, no unstructured dependency)."""
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        log.error("Failed to open PPTX file '%s': %s", file_path, exc, exc_info=True)
        raise
    texts = []
    for i, slide in enumerate(prs.slides, 1):   
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
        if parts:
            texts.append(Document(
                page_content="\n".join(parts),
                metadata={"source": file_path, "slide": i, "doc_type": "presentation"},
            ))
    log.debug("PPTX '%s': extracted %d slide(s) with text", file_path, len(texts))
    return texts


LOADER_MAP = {
    ".pdf": PyPDFLoader,
    ".docx": Docx2txtLoader,
}


def _enrich_metadata(chunks: list[Document], file_path: Path) -> list[Document]:
    """Add rich metadata tags to each chunk for hybrid retrieval filtering."""
    filename = file_path.name
    ext = file_path.suffix.lower()
    doc_type = {"pdf": "pdf", ".docx": "word", ".pptx": "presentation"}.get(ext, "document")
    total_chunks = len(chunks)
    for idx, chunk in enumerate(chunks):
        chunk.metadata.update({
            "filename": filename,
            "doc_type": doc_type,
            "chunk_index": idx,
            "total_chunks": total_chunks,
            # section position tag for prompting context
            "position": "start" if idx < total_chunks * 0.2 else (
                "end" if idx > total_chunks * 0.8 else "middle"
            ),
        })
        # Extract keyword hints from first line of chunk
        first_line = chunk.page_content.split("\n")[0].strip()[:120]
        chunk.metadata["heading_hint"] = first_line
    return chunks


def load_and_split(file_path: Path) -> list:
    """Load a document and split into chunks with rich metadata."""
    ext = file_path.suffix.lower()
    log.debug("Loading document '%s' (type=%s)", file_path.name, ext)

    if ext == ".pptx":
        documents = _load_pptx(str(file_path))
        chunks = SPLITTER.split_documents(documents)
        result = _enrich_metadata(chunks, file_path)
        log.info("Loaded '%s': %d chunk(s) from %d slide(s)", file_path.name, len(result), len(documents))
        return result

    loader_cls = LOADER_MAP.get(ext)
    if loader_cls is None:
        log.error("Unsupported file type '%s' for file '%s'", ext, file_path.name)
        raise ValueError(f"Unsupported file type: {ext}")
    try:
        loader = loader_cls(str(file_path))
        documents = loader.load()
    except Exception as exc:
        log.error("Failed to load '%s': %s", file_path.name, exc, exc_info=True)
        raise
    chunks = SPLITTER.split_documents(documents)
    result = _enrich_metadata(chunks, file_path)
    log.info("Loaded '%s': %d chunk(s) from %d page(s)", file_path.name, len(result), len(documents))
    return result

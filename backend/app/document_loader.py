"""Document loaders — extract text from PDF, DOCX, and PPTX files."""

from pathlib import Path

from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
)
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

SPLITTER = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)


def _load_pptx(file_path: str) -> list[Document]:
    """Load a PPTX file using python-pptx (fast, no unstructured dependency)."""
    prs = Presentation(file_path)
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
        if parts:
            texts.append(Document(
                page_content="\n".join(parts),
                metadata={"source": file_path, "slide": i, "doc_type": "presentation"},
            ))
    return texts


LOADER_MAP = {
    ".pdf": PyPDFLoader,
    ".docx": Docx2txtLoader,
}


def _enrich_metadata(chunks: list[Document], file_path: Path) -> list[Document]:
    """Add rich metadata tags to each chunk for hybrid retrieval filtering."""
    filename = file_path.name
    ext = file_path.suffix.lower()
    doc_type = {"pdf": "pdf", ".docx": "word", ".pptx": "presentation"}.get(ext, "document")
    total_chunks = len(chunks)
    for idx, chunk in enumerate(chunks):
        chunk.metadata.update({
            "filename": filename,
            "doc_type": doc_type,
            "chunk_index": idx,
            "total_chunks": total_chunks,
            # section position tag for prompting context
            "position": "start" if idx < total_chunks * 0.2 else (
                "end" if idx > total_chunks * 0.8 else "middle"
            ),
        })
        # Extract keyword hints from first line of chunk
        first_line = chunk.page_content.split("\n")[0].strip()[:120]
        chunk.metadata["heading_hint"] = first_line
    return chunks


def load_and_split(file_path: Path) -> list:
    """Load a document and split into chunks with rich metadata."""
    ext = file_path.suffix.lower()

    if ext == ".pptx":
        documents = _load_pptx(str(file_path))
        chunks = SPLITTER.split_documents(documents)
        return _enrich_metadata(chunks, file_path)

    loader_cls = LOADER_MAP.get(ext)
    if loader_cls is None:
        raise ValueError(f"Unsupported file type: {ext}")
    loader = loader_cls(str(file_path))
    documents = loader.load()
    chunks = SPLITTER.split_documents(documents)
    return _enrich_metadata(chunks, file_path)
